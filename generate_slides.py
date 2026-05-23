from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
FIGURES = ROOT / "figures"
OUTPUT = ROOT / "mad_system_presentation.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(20, 45, 80)
BLUE = RGBColor(35, 92, 150)
LIGHT_BLUE = RGBColor(230, 240, 250)
GREEN = RGBColor(35, 120, 95)
GRAY = RGBColor(90, 95, 105)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(25, 25, 25)
ORANGE = RGBColor(220, 130, 50)


def add_bg(slide, color=RGBColor(248, 250, 252)):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.65))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(27)
    run.font.bold = True
    run.font.color.rgb = DARK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.88), Inches(11.6), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Arial"
        sr.font.size = Pt(12)
        sr.font.color.rgb = GRAY
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.18), Inches(12.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE


def add_footer(slide, number):
    box = slide.shapes.add_textbox(Inches(11.95), Inches(7.08), Inches(0.9), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{number}/15"
    r.font.name = "Arial"
    r.font.size = Pt(9)
    r.font.color.rgb = GRAY


def add_bullets(slide, bullets, x, y, w, h, font_size=20, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.clear()
    for idx, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.space_after = Pt(8)
        p.font.name = "Arial"
        p.font.size = Pt(font_size if level == 0 else font_size - 2)
        p.font.color.rgb = color
    return box


def add_image(slide, image_name, x, y, w=None, h=None):
    path = FIGURES / image_name
    if not path.exists():
        return None
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if w is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(210, 220, 230)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    t = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.3), Inches(0.35))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Arial"
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = DARK
    b = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.52), Inches(w - 0.35), Inches(h - 0.62))
    bf = b.text_frame
    bf.word_wrap = True
    bf.clear()
    bp = bf.paragraphs[0]
    br = bp.add_run()
    br.text = body
    br.font.name = "Arial"
    br.font.size = Pt(12.5)
    br.font.color.rgb = BLACK


def add_notes(slide, notes):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes


def make_slide(title, number, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, subtitle)
    add_footer(slide, number)
    return slide


# 1. Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(242, 247, 252))
cover_title = slide.shapes.add_textbox(Inches(0.85), Inches(1.1), Inches(11.7), Inches(1.4))
tf = cover_title.text_frame
tf.clear()
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Nghiên cứu phát triển hệ thống phát hiện tin giả\ndựa trên phương pháp tranh biện đa tác tử"
r.font.name = "Arial"
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = DARK
sub = slide.shapes.add_textbox(Inches(1.4), Inches(2.8), Inches(10.5), Inches(0.6))
stf = sub.text_frame
stf.clear()
sp = stf.paragraphs[0]
sp.alignment = PP_ALIGN.CENTER
sr = sp.add_run()
sr.text = "MAD System for Fake News Detection"
sr.font.name = "Arial"
sr.font.size = Pt(22)
sr.font.color.rgb = BLUE
info = slide.shapes.add_textbox(Inches(2.2), Inches(4.2), Inches(8.9), Inches(1.1))
itf = info.text_frame
itf.clear()
for i, text in enumerate([
    "Sinh viên: Phí Đình Mạnh, Hoàng Ngọc Hưng, Trần Trung Kiên",
    "Giảng viên hướng dẫn: TS. Vương Thị Hồng",
    "Trường Đại học Công nghệ - Đại học Quốc gia Hà Nội",
]):
    pp = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
    pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run()
    rr.text = text
    rr.font.name = "Arial"
    rr.font.size = Pt(15)
    rr.font.color.rgb = BLACK
add_footer(slide, 1)
add_notes(slide, "Giới thiệu ngắn đề tài: hệ thống phát hiện tin giả dựa trên tranh biện đa tác tử, tập trung vào khả năng dùng bằng chứng và giải thích quyết định.")

# 2. Context
slide = make_slide("Bối cảnh và vấn đề", 2)
add_bullets(slide, [
    "Tin giả lan truyền nhanh trong xã hội số và gây ảnh hưởng tới y tế, chính trị, tài chính và niềm tin cộng đồng.",
    "Người dùng khó tự kiểm chứng vì thông tin thường thiếu nguồn, bị cắt ngữ cảnh hoặc đánh tráo khái niệm.",
    "Fact-checking hiện đại cần bằng chứng, nguồn gốc thông tin và lời giải thích có thể truy vết.",
    "LLM có tiềm năng hỗ trợ phân tích, nhưng cần cơ chế kiểm soát để tránh kết luận thiếu căn cứ.",
], 0.85, 1.55, 6.0, 4.9, 19)
add_card(slide, 7.35, 1.72, 4.9, 1.15, "Bài toán", "Đánh giá một bản tin hoặc claim văn bản là đáng tin hay sai lệch dựa trên bằng chứng.", BLUE)
add_card(slide, 7.35, 3.15, 4.9, 1.15, "Yêu cầu", "Không chỉ trả lời đúng/sai, mà phải chỉ ra lý do, nguồn bằng chứng và điểm quyết định.", GREEN)
add_card(slide, 7.35, 4.58, 4.9, 1.15, "Định hướng", "Kết hợp truy xuất bằng chứng với tranh biện đa tác tử để tăng kiểm tra chéo.", ORANGE)
add_notes(slide, "Nhấn mạnh đây không chỉ là bài toán phân loại văn bản. Với fake news detection, hệ thống phải chứng minh được vì sao kết luận là hợp lý.")

# 3. Single LLM limits
slide = make_slide("Vì sao LLM đơn lẻ chưa đủ?", 3)
add_card(slide, 0.8, 1.55, 3.8, 1.35, "Ảo giác thông tin", "LLM có thể tự tin khẳng định chi tiết không tồn tại trong nguồn bằng chứng.", ORANGE)
add_card(slide, 4.85, 1.55, 3.8, 1.35, "Kiến thức tĩnh", "Model không tự cập nhật sự kiện mới nếu không có cơ chế truy xuất bên ngoài.", BLUE)
add_card(slide, 8.9, 1.55, 3.8, 1.35, "Thiếu phản biện", "Một câu trả lời trực tiếp thường không tự soi lỗi thực thể, thời gian hoặc logic.", GREEN)
add_bullets(slide, [
    "Câu trả lời một chiều khó tạo audit trail cho người đọc hoặc hội đồng kiểm chứng.",
    "Kết quả dễ phụ thuộc vào prompt và cách diễn đạt claim đầu vào.",
    "Cần một cơ chế buộc lập luận phải đi qua đối soát bằng chứng và phản biện.",
], 1.15, 3.55, 10.8, 2.2, 21)
add_notes(slide, "Slide này dẫn vào lý do chọn MAD. Không phủ nhận năng lực LLM, mà chỉ ra rằng fact-checking cần thêm evidence và adversarial reasoning.")

# 4. Objectives
slide = make_slide("Mục tiêu và đóng góp của đồ án", 4)
add_bullets(slide, [
    "Xây dựng hệ thống kiểm chứng tin giả cho dữ liệu văn bản bằng Multi-Agent Debate.",
    "Thiết kế workflow có state chung để lưu knowledge_base, source_scores, debate_history, claims_registry và verdict.",
    "Hỗ trợ hai chế độ: Search Mode cho bối cảnh mở và Non-search Mode cho benchmark có evidence sẵn.",
    "Tạo cơ chế lập luận có truy vết thông qua claim ID D*/C* và lịch sử tranh luận nhiều vòng.",
    "Đánh giá thực nghiệm Base LLM so với MAD trên FEVER binary SUPPORTS/REFUTES.",
], 0.85, 1.45, 11.7, 4.8, 20)
add_notes(slide, "Trình bày mục tiêu theo hướng hệ thống: không chỉ tạo model, mà xây dựng pipeline có bằng chứng, trạng thái, tranh luận và đánh giá.")

# 5. MAD idea
slide = make_slide("Ý tưởng Multi-Agent Debate", 5, "Mô hình phiên tòa số cho kiểm chứng thông tin")
add_image(slide, "mad_debate_court.png", 0.7, 1.45, w=6.1)
add_bullets(slide, [
    "Defender bảo vệ khả năng bản tin đúng.",
    "Challenger tìm điểm sai, thiếu căn cứ hoặc đánh tráo khái niệm.",
    "Source Scorer đánh giá độ tin cậy của nguồn.",
    "Judge tổng hợp bằng chứng và lịch sử tranh luận để đưa ra truth_score.",
], 7.15, 1.65, 5.45, 4.6, 20)
add_notes(slide, "Giải thích hình như một phiên tòa: hai bên tranh luận, nguồn đóng vai trò bằng chứng, Judge ra phán quyết cuối cùng.")

# 6. System overview
slide = make_slide("Tổng quan kiến trúc hệ thống", 6)
add_image(slide, "mad_system_overview.png", 0.55, 1.35, w=12.2)
add_notes(slide, "Đây là slide trung tâm. Đi từ input, truy xuất bằng chứng, MADState, debate, Judge đến output. Nhấn mạnh state giúp dữ liệu không đi một chiều mà được tái sử dụng qua nhiều vòng.")

# 7. Components
slide = make_slide("Các thành phần chính", 7)
components = [
    ("Defender", "Tạo D1, D2... để bảo vệ bản tin dựa trên evidence."),
    ("Challenger", "Tạo C1, C2... để phản biện lỗi thực thể, thời gian, logic."),
    ("Search Module", "Lập kế hoạch truy vấn theo từng phía, Tavily trước, Wikipedia fallback."),
    ("Source Scorer", "Chấm trust_score cho nguồn để Judge cân nhắc."),
    ("Claims Registry", "Lưu lịch sử ASSERT, REBUT, DEFEND theo từng claim."),
    ("Judge", "Tổng hợp evidence, source_scores và debate_history thành verdict."),
]
for idx, (title, body) in enumerate(components):
    row = idx // 3
    col = idx % 3
    add_card(slide, 0.7 + col * 4.18, 1.55 + row * 2.05, 3.75, 1.55, title, body, [BLUE, GREEN, ORANGE][col])
add_notes(slide, "Slide này giúp hội đồng nắm được vai trò từng module trước khi đi vào state và workflow.")

# 8. State and claims
slide = make_slide("MADState và Claims Registry", 8)
add_image(slide, "mad_state_claims_registry.png", 0.65, 1.38, w=7.1)
add_bullets(slide, [
    "MADState là bộ nhớ làm việc chung của toàn bộ workflow.",
    "knowledge_base và source_scores lưu bằng chứng và độ tin cậy nguồn.",
    "debate_history lưu toàn bộ diễn biến qua các vòng.",
    "claims_registry giữ lịch sử D*/C*, giúp REBUT/DEFEND nhắm đúng claim.",
], 8.05, 1.65, 4.55, 4.8, 18)
add_notes(slide, "Nhấn mạnh đây là điểm khác biệt quan trọng so với hỏi LLM trực tiếp: hệ thống có trạng thái, có ký ức tranh luận và có cấu trúc claim.")

# 9. Search workflow
slide = make_slide("Luồng xử lý Search Mode", 9)
add_image(slide, "search_workflow.png", 0.55, 1.28, w=12.15)
add_notes(slide, "Giải thích Search Mode: prepare_round tạo kế hoạch, hai phía tìm kiếm, score_sources, defender, challenger, save_round, lặp đến max_rounds rồi judge.")

# 10. Non-search workflow
slide = make_slide("Luồng xử lý Non-search Mode", 10)
add_image(slide, "non_search_workflow.png", 0.8, 1.38, w=7.0)
add_bullets(slide, [
    "Dùng khi context/evidence đã có sẵn, ví dụ FEVER.",
    "initial_context được nạp thành nguồn [S1] với trust score 1.0.",
    "Không gọi Tavily/Wikipedia, giúp đánh giá ổn định và kiểm soát input.",
    "Vẫn giữ cấu trúc Defender - Challenger - Save Round - Judge.",
], 8.15, 1.72, 4.35, 4.45, 18)
add_notes(slide, "Slide này nối sang phần thực nghiệm. FEVER cung cấp claim và evidence nên không cần search web.")

# 11. Implementation stack
slide = make_slide("Cài đặt và công nghệ sử dụng", 11)
techs = [
    ("Python", "Ngôn ngữ triển khai chính."),
    ("LangGraph", "Điều phối workflow dạng StateGraph có vòng lặp."),
    ("LangChain ChatOpenAI", "Kết nối LLM qua API tương thích OpenAI."),
    ("NineRouter", "Provider/model cấu hình qua biến môi trường."),
    ("Tavily & Wikipedia", "Truy xuất bằng chứng trong Search Mode."),
    ("FEVER", "Benchmark SUPPORTS/REFUTES cho đánh giá thực nghiệm."),
]
for idx, (title, body) in enumerate(techs):
    row = idx // 2
    col = idx % 2
    add_card(slide, 0.9 + col * 6.05, 1.42 + row * 1.55, 5.55, 1.1, title, body, BLUE if col == 0 else GREEN)
add_notes(slide, "Không cần đọc hết công nghệ. Tập trung vào LangGraph cho workflow, LLM qua NineRouter và Tavily/Wikipedia cho evidence.")

# 12. FEVER pipeline
slide = make_slide("Quy trình thực nghiệm FEVER", 12)
add_image(slide, "fever_evaluation_pipeline.png", 0.6, 1.35, w=12.0)
add_notes(slide, "Giải thích quy trình: chuẩn bị dữ liệu FEVER, lấy claim/evidence, chạy Base LLM và MAD non-search, chuyển truth_score về binary rồi tính metric.")

# 13. Results
slide = make_slide("Kết quả thực nghiệm", 13)
add_image(slide, "fever_accuracy_comparison.png", 0.65, 1.35, w=7.0)
add_bullets(slide, [
    "Llama 3.3 70B: 92.5% -> 97.5% (+5.0%).",
    "Gemma 4-31B: 82.5% -> 90.0% (+7.5%).",
    "GPT-OSS 120B: 85.0% -> 92.5% (+7.5%).",
    "Gemini 3.1 Flash Lite: 80.0% -> 85.0% (+5.0%).",
    "MAD cải thiện ổn định trên các model trong thiết lập thử nghiệm FEVER.",
], 8.0, 1.5, 4.7, 5.0, 17)
add_notes(slide, "Nhấn mạnh kết quả chính: MAD không chỉ tăng accuracy mà còn tạo lời giải thích và lịch sử tranh luận để truy vết.")

# 14. Analysis and limitations
slide = make_slide("Phân tích, hạn chế và hướng phát triển", 14)
add_card(slide, 0.8, 1.45, 3.8, 4.75, "Nhận xét", "MAD thận trọng hơn hỏi LLM trực tiếp, giảm xu hướng kết luận thiếu căn cứ và giúp giải thích quyết định qua evidence + debate history.", GREEN)
add_card(slide, 4.85, 1.45, 3.8, 4.75, "Hạn chế", "Chi phí thời gian/token cao hơn; phụ thuộc chất lượng LLM và nguồn; structured output vẫn có thể lỗi với model nhỏ.", ORANGE)
add_card(slide, 8.9, 1.45, 3.8, 4.75, "Phát triển", "Đồng bộ UI, dùng schema/tool calling, mở rộng GossipCop/TruthfulQA, model routing, multilingual và multimodal verification.", BLUE)
add_notes(slide, "Slide này chuẩn bị cho câu hỏi phản biện. Nên chủ động nói trade-off: tăng khả năng kiểm chứng nhưng tốn tài nguyên hơn.")

# 15. Conclusion
slide = make_slide("Kết luận", 15)
add_bullets(slide, [
    "Đồ án đã xây dựng hệ thống phát hiện tin giả dựa trên tranh biện đa tác tử.",
    "Hệ thống kết hợp truy xuất bằng chứng, chấm điểm nguồn, quản lý nhận định và phán quyết cuối cùng.",
    "Thực nghiệm FEVER cho thấy MAD cải thiện accuracy so với Base LLM trong các lần chạy đã ghi nhận.",
    "Giá trị chính của hệ thống nằm ở khả năng giải thích, truy vết và kiểm tra chéo lập luận.",
], 1.2, 1.55, 10.9, 3.5, 22)
q = slide.shapes.add_textbox(Inches(4.2), Inches(5.75), Inches(4.9), Inches(0.7))
qtf = q.text_frame
qtf.clear()
qp = qtf.paragraphs[0]
qp.alignment = PP_ALIGN.CENTER
qr = qp.add_run()
qr.text = "Q&A"
qr.font.name = "Arial"
qr.font.size = Pt(36)
qr.font.bold = True
qr.font.color.rgb = BLUE
add_notes(slide, "Kết luận ngắn, sau đó chuyển sang hỏi đáp. Nếu có thời gian, nhắc lại slide tổng quan kiến trúc khi trả lời câu hỏi.")

prs.save(OUTPUT)
print(f"Saved {OUTPUT}")
